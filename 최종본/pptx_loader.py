# -*- coding: utf-8 -*-
"""PPTX(파워포인트) 로더.

app.py 가 pdf_loader.process_pdf 의 반환값(result)에서 사용하는 필드와
동일한 형태의 결과 객체를 만들어, PPT 도 PDF 와 똑같은 파이프라인을 타게 한다.

사용하는 필드(app.py / summarize_page / output_page 기준):
    result.file_name
    result.page_count
    result.full_text
    result.ai_mode.mode
    result.ai_mode.average_chars_per_page
    result.ai_mode.reason
    result.page_images            (PPT는 슬라이드 렌더 이미지가 없으므로 [])
    result.pages[i].page_number
    result.pages[i].cleaned_text

의존성: python-pptx  (pip install python-pptx)
"""
from __future__ import annotations
from dataclasses import dataclass, field
from typing import List
import io


@dataclass
class AiMode:
    mode: str
    average_chars_per_page: float
    reason: str


@dataclass
class SlidePage:
    page_number: int
    cleaned_text: str

    @property
    def text(self) -> str:        # 호환용 별칭
        return self.cleaned_text


@dataclass
class PptxResult:
    file_name: str
    page_count: int
    full_text: str
    ai_mode: AiMode
    pages: List[SlidePage]
    page_images: list = field(default_factory=list)


def _iter_shapes(shapes):
    """그룹 도형까지 재귀적으로 순회."""
    for shape in shapes:
        # 그룹 도형(6 = MSO_SHAPE_TYPE.GROUP)
        if getattr(shape, "shape_type", None) == 6 and hasattr(shape, "shapes"):
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _shape_text(shape) -> List[str]:
    parts: List[str] = []
    # 텍스트 프레임
    if getattr(shape, "has_text_frame", False):
        txt = (shape.text_frame.text or "").strip()
        if txt:
            parts.append(txt)
    # 표
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [(c.text or "").strip() for c in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return parts


def extract_slides(prs) -> List[str]:
    pages: List[str] = []
    for slide in prs.slides:
        parts: List[str] = []
        for shape in _iter_shapes(slide.shapes):
            parts.extend(_shape_text(shape))
        # 발표자 노트
        try:
            if slide.has_notes_slide:
                note = (slide.notes_slide.notes_text_frame.text or "").strip()
                if note:
                    parts.append("[발표자 노트] " + note)
        except Exception:
            pass
        pages.append("\n".join(parts).strip())
    return pages


def process_pptx(file_bytes: bytes, file_name: str) -> PptxResult:
    """PPTX 바이트 → app.py 호환 결과 객체."""
    try:
        from pptx import Presentation
    except ImportError as e:
        raise RuntimeError(
            "python-pptx 가 설치되어 있지 않습니다. "
            "`pip install python-pptx` 후 다시 시도하세요."
        ) from e

    prs = Presentation(io.BytesIO(file_bytes))
    slide_texts = extract_slides(prs)

    pages = [SlidePage(page_number=i, cleaned_text=t)
             for i, t in enumerate(slide_texts, 1)]
    page_count = len(pages) or 1
    full_text = "\n\n".join(
        f"[슬라이드 {p.page_number}]\n{p.cleaned_text}" for p in pages if p.cleaned_text
    )
    avg = (len(full_text) / page_count) if page_count else 0.0

    ai_mode = AiMode(
        mode="PPT 텍스트",
        average_chars_per_page=avg,
        reason=(
            f"PowerPoint 파일에서 슬라이드 {page_count}장의 텍스트(본문·표·발표자 노트)를 "
            f"직접 추출했습니다. 슬라이드 렌더 이미지는 제공되지 않으며, 추출된 텍스트로 요약을 진행합니다."
        ),
    )
    return PptxResult(
        file_name=file_name,
        page_count=page_count,
        full_text=full_text,
        ai_mode=ai_mode,
        pages=pages,
        page_images=[],
    )
